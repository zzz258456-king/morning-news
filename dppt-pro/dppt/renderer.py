"""
DPPT PPTX 渲染器
把 DpptDocument 渲染为 PowerPoint 文件。
"""

from __future__ import annotations

import io
import os
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from .animation import add_entrance_animation
from .icons import render_icon
from .models import (
    Chart,
    DpptDocument,
    Fill,
    Icon,
    Image,
    Page,
    Shape,
    Table,
    TableCell,
    Text,
    Theme,
)
from .theme import apply_theme


# 16:9 默认尺寸 (inches)
DEFAULT_SLIDE_WIDTH = Inches(13.333)
DEFAULT_SLIDE_HEIGHT = Inches(7.5)


def _to_inches(value: float, unit: str = "pt") -> float:
    """把 pt/px/mm 转为 inch，便于 python-pptx 使用。"""
    if unit == "inch":
        return value
    if unit == "pt":
        return value / 72.0
    if unit == "px":
        return value / 96.0
    if unit == "mm":
        return value / 25.4
    return value


def _pt_to_inches(value: Optional[Any]) -> Optional[float]:
    if value is None:
        return None
    if hasattr(value, "value"):
        return _to_inches(value.value, value.unit)
    return _to_inches(float(value), "pt")


def _hex_to_rgb(color: str) -> Optional[Tuple[int, int, int]]:
    if not color:
        return None
    color = color.lstrip("#")
    if len(color) == 3:
        color = "".join([c * 2 for c in color])
    if len(color) == 6:
        try:
            return tuple(int(color[i : i + 2], 16) for i in (0, 2, 4))
        except ValueError:
            return None
    return None


def _apply_fill(shape, fill: Optional[Fill]) -> None:
    if fill is None:
        return
    if fill.type == "solid" and fill.color:
        rgb = _hex_to_rgb(fill.color)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
    elif fill.type == "gradient" and fill.gradient:
        # python-pptx 对渐变支持有限，回退为第一种 stop 颜色
        stops = fill.gradient.get("stops", [])
        if stops:
            first = stops[0]
            color = first.get("color") if isinstance(first, dict) else first
            rgb = _hex_to_rgb(color)
            if rgb:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)
    elif fill.type == "none":
        shape.fill.background()


def _apply_line(shape, border) -> None:
    if border is None or border.width <= 0:
        shape.line.fill.background()
        return
    rgb = _hex_to_rgb(border.color)
    if rgb:
        shape.line.color.rgb = RGBColor(*rgb)
    shape.line.width = Pt(border.width)


def _apply_text_style(run, style) -> None:
    if style.font_family:
        run.font.name = style.font_family
    if style.font_size is not None:
        inches = _pt_to_inches(style.font_size)
        if inches is not None:
            run.font.size = Pt(inches * 72)
    if style.color:
        rgb = _hex_to_rgb(style.color)
        if rgb:
            run.font.color.rgb = RGBColor(*rgb)
    run.font.bold = style.bold
    run.font.italic = style.italic
    run.font.underline = style.underline


def _align_to_pp(align: str):
    mapping = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    return mapping.get(align, PP_ALIGN.LEFT)


def _valign_to_anchor(valign: str):
    mapping = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    return mapping.get(valign, MSO_ANCHOR.TOP)


def _render_text(slide, text_elem: Text):
    b = text_elem.bounds
    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    width = Inches(b.get("width", 1))
    height = Inches(b.get("height", 0.5))

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = text_elem.wrap
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if text_elem.auto_size else MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = _valign_to_anchor(text_elem.style.valign)

    p = tf.paragraphs[0]
    p.alignment = _align_to_pp(text_elem.style.align)
    run = p.add_run()
    run.text = text_elem.text
    _apply_text_style(run, text_elem.style)
    return shape


def _render_shape(slide, shape_elem: Shape):
    b = shape_elem.bounds
    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    width = Inches(b.get("width", 1))
    height = Inches(b.get("height", 1))

    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "line": MSO_SHAPE.LINE_INVERSE,
    }
    mso = mapping.get(shape_elem.shape_type, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(mso, left, top, width, height)
    _apply_fill(shape, shape_elem.fill)
    _apply_line(shape, shape_elem.border)
    return shape


def _resolve_image_path(src: str, doc: DpptDocument) -> Optional[str]:
    if os.path.isabs(src) and os.path.exists(src):
        return src
    base = doc.resources.get("__base__", "")
    candidate = os.path.join(base, src)
    if os.path.exists(candidate):
        return candidate
    return None


def _render_image(slide, img_elem: Image, doc: DpptDocument):
    b = img_elem.bounds
    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    width = Inches(b.get("width", 1))
    height = Inches(b.get("height", 1))

    path = _resolve_image_path(img_elem.src, doc)
    if path:
        return slide.shapes.add_picture(path, left, top, width, height)
    else:
        # 占位矩形
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        _apply_fill(shape, Fill(type="solid", color="#E0E0E0"))
        return shape


def _render_icon(slide, icon_elem: Icon, doc: DpptDocument):
    b = icon_elem.bounds
    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    width = Inches(b.get("width", 0.3))
    height = Inches(b.get("height", 0.3))

    color = icon_elem.color or (icon_elem.fill.color if icon_elem.fill else None)
    size_px = int(max(width, height) * 96 * 2)  # 2x 分辨率保证清晰

    png_bytes = render_icon(
        icon_name=icon_elem.icon_name,
        svg_path=icon_elem.svg_path,
        svg_file=icon_elem.src,
        color=color,
        size=max(size_px, 32),
    )
    if png_bytes:
        buf = io.BytesIO(png_bytes)
        return slide.shapes.add_picture(buf, left, top, width, height)
    else:
        # fallback：占位矩形
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        _apply_fill(shape, Fill(type="solid", color=color or "#666666"))
        return shape


def _render_table(slide, table_elem: Table):
    b = table_elem.bounds
    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    width = Inches(b.get("width", 4))
    height = Inches(b.get("height", 2))

    rows = len(table_elem.rows)
    cols = max((len(row) for row in table_elem.rows), default=0)
    if rows == 0 or cols == 0:
        return None

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table

    for r_idx, row in enumerate(table_elem.rows):
        for c_idx, cell in enumerate(row):
            if c_idx >= cols:
                break
            pcell = tbl.cell(r_idx, c_idx)
            pcell.text = cell.text
            if cell.fill and cell.fill.color:
                rgb = _hex_to_rgb(cell.fill.color)
                if rgb:
                    pcell.fill.solid()
                    pcell.fill.fore_color.rgb = RGBColor(*rgb)
            # 简单文本样式
            for paragraph in pcell.text_frame.paragraphs:
                paragraph.alignment = _align_to_pp(cell.style.align)
                for run in paragraph.runs:
                    _apply_text_style(run, cell.style)
    return shape


def _render_chart(slide, chart_elem: Chart) -> None:
    # MVP：使用 matplotlib 生成图表图片再插入
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import font_manager

        # 尝试设置中文字体
        if not hasattr(_render_chart, "_chinese_font"):
            candidates = ["Microsoft YaHei", "SimHei", "WenQuanYi Micro Hei", "Noto Sans CJK SC"]
            found = None
            for name in candidates:
                try:
                    fm = font_manager.FontProperties(family=name)
                    if fm.get_name() != "DejaVu Sans":
                        found = name
                        break
                except Exception:
                    continue
            _render_chart._chinese_font = found
        if _render_chart._chinese_font:
            plt.rcParams["font.family"] = _render_chart._chinese_font
            plt.rcParams["axes.unicode_minus"] = False
    except Exception:
        return

    b = chart_elem.bounds
    width = b.get("width", 4)
    height = b.get("height", 3)
    dpi = 150

    fig, ax = plt.subplots(figsize=(width, height), dpi=dpi)
    labels = chart_elem.labels
    series = chart_elem.series

    if chart_elem.chart_type == "bar":
        for idx, s in enumerate(series):
            values = s.get("data", [])
            ax.bar(labels, values, label=s.get("name", f"Series {idx}"), color=chart_elem.colors[idx] if idx < len(chart_elem.colors) else None)
        if chart_elem.show_legend:
            ax.legend()
    elif chart_elem.chart_type == "line":
        for idx, s in enumerate(series):
            values = s.get("data", [])
            ax.plot(labels, values, label=s.get("name", f"Series {idx}"), color=chart_elem.colors[idx] if idx < len(chart_elem.colors) else None)
        if chart_elem.show_legend:
            ax.legend()
    elif chart_elem.chart_type == "pie" and series:
        values = series[0].get("data", [])
        ax.pie(values, labels=labels, colors=chart_elem.colors[: len(values)] if chart_elem.colors else None)

    if chart_elem.title:
        ax.set_title(chart_elem.title)
    if chart_elem.show_grid:
        ax.grid(True, linestyle="--", alpha=0.3)

    buf = io.BytesIO()
    plt.tight_layout()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    left = Inches(b.get("x", 0))
    top = Inches(b.get("y", 0))
    return slide.shapes.add_picture(buf, left, top, Inches(width), Inches(height))


def _render_element(slide, elem, doc: DpptDocument):
    if not elem.visible:
        return None
    if elem.type == "text":
        return _render_text(slide, elem)
    elif elem.type == "shape":
        return _render_shape(slide, elem)
    elif elem.type == "image":
        return _render_image(slide, elem, doc)
    elif elem.type == "icon":
        return _render_icon(slide, elem, doc)
    elif elem.type == "table":
        return _render_table(slide, elem)
    elif elem.type == "chart":
        return _render_chart(slide, elem)
    return None


def _set_slide_background(slide, fill: Optional[Fill]) -> None:
    if fill is None:
        return
    if fill.type == "solid" and fill.color:
        rgb = _hex_to_rgb(fill.color)
        if rgb:
            background = slide.background
            fill_format = background.fill
            fill_format.solid()
            fill_format.fore_color.rgb = RGBColor(*rgb)


class DpptRenderer:
    """DPPT PPTX 渲染器。"""

    def __init__(self, document: DpptDocument, enable_animation: bool = True):
        self.document = document
        self.enable_animation = enable_animation
        apply_theme(document)

    def render(self, output_path: str, master_path: Optional[str] = None) -> None:
        """渲染为 PPTX 文件。"""
        if master_path and os.path.exists(master_path):
            prs = Presentation(master_path)
        else:
            prs = Presentation()
            prs.slide_width = DEFAULT_SLIDE_WIDTH
            prs.slide_height = DEFAULT_SLIDE_HEIGHT

        # 使用空白布局
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        for page in self.document.pages:
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_background(slide, page.background)
            animated_shapes = []
            for elem in page.elements:
                shape = _render_element(slide, elem, self.document)
                if shape and self.enable_animation and elem.type in ("text", "image", "icon", "table", "chart"):
                    animated_shapes.append(shape)
            if self.enable_animation and animated_shapes:
                for idx, shape in enumerate(animated_shapes):
                    add_entrance_animation(
                        slide,
                        shape,
                        effect="fade",
                        duration_ms=400,
                        trigger="onClick",
                        delay_ms=0,
                    )

        prs.save(output_path)

    def render_to_bytes(self, master_path: Optional[str] = None) -> bytes:
        buf = io.BytesIO()
        # python-pptx 不直接支持写入 BytesIO，先写临时文件
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            self.render(tmp_path, master_path=master_path)
            with open(tmp_path, "rb") as f:
                return f.read()
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
