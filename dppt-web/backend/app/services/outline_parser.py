"""
大纲解析服务
支持 txt / md / docx / pdf / pptx 等格式。
"""

from pathlib import Path


def parse_outline(file_path: Path) -> list:
    """解析大纲文件为页面列表"""
    suffix = file_path.suffix.lower()

    if suffix in [".txt", ".md"]:
        text = file_path.read_text(encoding="utf-8")
    elif suffix == ".docx":
        text = _parse_docx(file_path)
    elif suffix == ".pdf":
        text = _parse_pdf(file_path)
    elif suffix in [".pptx", ".ppt"]:
        text = _parse_ppt(file_path)
    else:
        # 兜底：尝试按文本读取
        text = file_path.read_text(encoding="utf-8")

    return _text_to_pages(text)


def _text_to_pages(text: str) -> list:
    pages = []
    idx = 1
    for line in text.strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        title = line.lstrip("#").strip().lstrip("*").strip().lstrip("-").strip()
        if title:
            pages.append({
                "id": f"page_{idx}",
                "title": title,
                "content": "",
                "notes": "",
            })
            idx += 1
    return pages


def _parse_docx(file_path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        raise ValueError("缺少 python-docx，无法解析 .docx 文件")


def _parse_pdf(file_path: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(file_path))
        return "\n".join(page.get_text() for page in doc)
    except ImportError:
        raise ValueError("缺少 PyMuPDF，无法解析 .pdf 文件")


def _parse_ppt(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except ImportError:
        raise ValueError("缺少 python-pptx，无法解析 .pptx 文件")
