import subprocess
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from app.models import GenerateRequest
from app.store import load_project, save_project, get_project_path

router = APIRouter()


def build_prompt(config: dict) -> str:
    slides_desc = "\n".join(
        f"- {s['page_id']}: {s['title']}\n  内容：{s.get('body', '')}"
        for s in config.get("slides", [])
    )
    template = config.get("template", {})
    return f"""请使用 ppt-from-outline-to-delivery 技能，根据以下配置生成一份 PPT：

主题：{config.get('title', '未命名')}
模板：{template.get('name', '默认')}（配色：{', '.join(template.get('colors', []))}，比例：{template.get('layout', '16:9')}）

页面内容：
{slides_desc}

要求：
1. 每页有明确的标题和视觉元素
2. 图片不裁切关键内容
3. 配色统一
4. 输出文件保存到：{get_project_path(config['id']) / 'output.pptx'}
"""


@router.post("/{project_id}/generate")
def generate_ppt(project_id: str, request: GenerateRequest):
    project = load_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    project_dir = get_project_path(project_id)
    project_dir.mkdir(parents=True, exist_ok=True)
    output_file = project_dir / "output.pptx"

    # 同步项目状态
    config_dict = request.config.model_dump()
    project["title"] = config_dict.get("title", project.get("title", ""))
    project["template"] = config_dict.get("template")
    project["slides"] = config_dict.get("slides", [])
    save_project(project_id, project)

    # 如果 slides 为空，根据 outline 生成默认 slides
    if not project["slides"] and project.get("outline"):
        project["slides"] = [
            {
                "page_id": p["id"],
                "title": p["title"],
                "body": p.get("content", ""),
                "image": None,
                "layout": "default",
            }
            for p in project["outline"]
        ]
        save_project(project_id, project)

    try:
        # TODO: 实际调用 Claude Code CLI
        # prompt = build_prompt({**config_dict, "id": project_id})
        # subprocess.run(
        #     ["claude", prompt],
        #     cwd=str(project_dir),
        #     timeout=300,
        #     check=True,
        # )

        # MVP 阶段：使用 python-pptx 生成一个基础 PPTX 验证链路
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = 12192000
        prs.slide_height = 6858000

        # 封面
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        txBox.text_frame.paragraphs[0].text = project.get("title", "DPPT 生成示例")
        txBox.text_frame.paragraphs[0].font.size = Pt(36)

        # 内容页
        for s in project.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            txBox.text_frame.paragraphs[0].text = s.get("title", "")
            txBox.text_frame.paragraphs[0].font.size = Pt(28)

            body = s.get("body", "")
            if body:
                txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
                txBox2.text_frame.paragraphs[0].text = body
                txBox2.text_frame.paragraphs[0].font.size = Pt(18)

        prs.save(str(output_file))

        project["output_path"] = str(output_file)
        save_project(project_id, project)

        return {
            "project_id": project_id,
            "status": "success",
            "output_path": str(output_file),
            "message": "PPT 已生成（MVP 使用占位内容）",
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成失败：{str(e)}")


@router.get("/{project_id}/download")
def download_ppt(project_id: str):
    project = load_project(project_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    output_file = Path(project.get("output_path", get_project_path(project_id) / "output.pptx"))
    if not output_file.exists():
        raise HTTPException(status_code=404, detail="文件不存在")

    return FileResponse(
        path=str(output_file),
        filename=f"{project.get('title', 'output')}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
